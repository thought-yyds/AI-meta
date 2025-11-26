"""FastMCP service that exposes the file_parser tool."""

from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Set

from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel, Field, ValidationError

try:
    from dotenv import load_dotenv
except ImportError:  # pragma: no cover - optional dependency
    load_dotenv = None


if load_dotenv is not None:
    load_dotenv(override=False)

BASE_WORKDIR = Path(os.getenv("PERSONAL_AGENT_WORKDIR", ".")).resolve()
TEXT_EXTENSIONS: Set[str] = {".txt", ".md", ".markdown"}

app = FastMCP("file-parser")


@dataclass
class MCPResponse:
    ok: bool
    data: Optional[Dict[str, object]] = None
    error: Optional[str] = None

    def to_dict(self) -> Dict[str, object]:
        result: Dict[str, object] = {"ok": self.ok}
        if self.data is not None:
            result["data"] = self.data
        if self.error:
            result["error"] = self.error
        return result


class FileParserSchema(BaseModel):
    path: str = Field(..., description="Path to the document to parse (relative or absolute).")
    max_sections: Optional[int] = Field(
        None,
        description="Optional limit for slides (pptx) or pages (pdf/docx). For text files, limits the number of lines.",
    )

    @property
    def resolved_path(self) -> Path:
        raw_path = Path(self.path)
        if raw_path.is_absolute():
            return raw_path.resolve()
        return (BASE_WORKDIR / raw_path).resolve()


@app.tool(name="file_parser")
def file_parser(path: str, max_sections: Optional[int] = None) -> Dict[str, object]:
    """Parse PPTX/PDF/DOCX/TXT/MD files and return structured content."""
    try:
        params = FileParserSchema(path=path, max_sections=max_sections)
        if params.max_sections is not None and params.max_sections <= 0:
            raise ValueError("max_sections must be a positive integer.")

        file_path = params.resolved_path
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()
        if suffix == ".pptx":
            payload = _parse_pptx(file_path, params.max_sections)
        elif suffix == ".pdf":
            payload = _parse_pdf(file_path, params.max_sections)
        elif suffix == ".docx":
            payload = _parse_docx(file_path, params.max_sections)
        elif suffix in TEXT_EXTENSIONS:
            payload = _parse_text(file_path, params.max_sections)
        else:
            supported = ", ".join(sorted({".pptx", ".pdf", ".docx", *TEXT_EXTENSIONS}))
            raise ValueError(f"Unsupported file extension '{suffix}'. Supported extensions: {supported}.")

        return MCPResponse(ok=True, data={"file": str(file_path), **payload}).to_dict()
    except ValidationError as exc:
        return MCPResponse(ok=False, error=str(exc)).to_dict()
    except Exception as exc:  # noqa: BLE001
        return MCPResponse(ok=False, error=str(exc)).to_dict()


def _parse_pptx(file_path: Path, max_sections: Optional[int]) -> Dict[str, object]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # noqa: BLE001
        raise RuntimeError("Missing dependency 'python-pptx'. Install it with 'pip install python-pptx'.") from exc

    presentation = Presentation(str(file_path))
    slides_data: List[Dict[str, object]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        if max_sections is not None and idx > max_sections:
            break
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {idx}"
        bullets: List[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if text:
                bullets.append(text)
        slides_data.append({"index": idx, "title": title.strip(), "content": bullets})
    return {"type": "pptx", "slide_count": len(slides_data), "slides": slides_data}


def _parse_pdf(file_path: Path, max_sections: Optional[int]) -> Dict[str, object]:
    try:
        from PyPDF2 import PdfReader  # type: ignore
    except ImportError as exc:  # noqa: BLE001
        raise RuntimeError("Missing dependency 'PyPDF2'. Install it with 'pip install PyPDF2'.") from exc

    reader = PdfReader(str(file_path))
    pages = []
    for idx, page in enumerate(reader.pages, start=1):
        if max_sections is not None and idx > max_sections:
            break
        text = page.extract_text() or ""
        pages.append({"index": idx, "content": text.strip()})
    return {"type": "pdf", "page_count": len(reader.pages), "pages": pages}


def _parse_docx(file_path: Path, max_sections: Optional[int]) -> Dict[str, object]:
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:  # noqa: BLE001
        raise RuntimeError("Missing dependency 'python-docx'. Install it with 'pip install python-docx'.") from exc

    document = Document(str(file_path))
    paragraphs: List[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)
        if max_sections is not None and len(paragraphs) >= max_sections:
            break

    return {"type": "docx", "paragraphs": paragraphs, "paragraph_count": len(document.paragraphs)}


def _parse_text(file_path: Path, max_sections: Optional[int]) -> Dict[str, object]:
    try:
        content = file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = file_path.read_text(encoding="utf-8", errors="ignore")

    lines = [line.strip() for line in content.splitlines()]
    if max_sections is not None:
        lines = lines[:max_sections]
    return {"type": "text", "lines": lines, "line_count": len(content.splitlines())}


if __name__ == "__main__":
    app.run()

